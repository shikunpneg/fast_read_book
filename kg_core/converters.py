"""
统一文件转换器：支持 PDF/EPUB/DOCX/PPTX/XLSX/HTML/MD/图片 等格式
- 文本型格式：直接提取
- 图像型 PDF/图片：调用 OCR
"""
import os
import re


def read_text(filepath: str) -> str:
    """读取文本文件，自动检测编码"""
    encodings = ['utf-8', 'utf-8-sig', 'gbk', 'gb2312', 'gb18030', 'big5', 'latin-1']
    for enc in encodings:
        try:
            with open(filepath, encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    with open(filepath, encoding='utf-8', errors='replace') as f:
        return f.read()


def extract_pdf(filepath: str) -> str:
    """从 PDF 提取文本（PyMuPDF）"""
    try:
        import fitz
    except ImportError:
        raise RuntimeError("需要安装 PyMuPDF: pip install PyMuPDF")
    doc = fitz.open(filepath)
    parts = [page.get_text() for page in doc]
    doc.close()
    text = '\n\n'.join(parts)
    # 如果文本极少（< 100 字/页），判定为图像型 PDF
    if len(text) < len(parts) * 100 and len(parts) > 1:
        raise _EmptyPDFError(f"PDF 可能是图像型（{len(parts)} 页仅 {len(text)} 字符），需要 OCR")
    if not text.strip():
        raise _EmptyPDFError("PDF 无可提取文本（图像型）")
    return text


def extract_epub(filepath: str) -> str:
    """从 EPUB 提取文本"""
    try:
        from ebooklib import epub
        from lxml import etree
    except ImportError:
        raise RuntimeError("需要安装 ebooklib: pip install ebooklib lxml")
    book = epub.read_epub(filepath)
    parts = []
    for item in book.get_items():
        if item.get_type() == 9:  # ITEM_DOCUMENT
            content = item.get_content().decode('utf-8', errors='replace')
            try:
                root = etree.fromstring(content.encode('utf-8'))
                for elem in root.iter():
                    if elem.tag in ('p', 'h1', 'h2', 'h3', 'h4', 'h5', 'li', 'div'):
                        text = ''.join(elem.itertext()).strip()
                        if text:
                            parts.append(text)
            except Exception:
                parts.append(content)
    return '\n\n'.join(parts)


def extract_docx(filepath: str) -> str:
    """从 Word DOCX 提取文本"""
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("需要安装 python-docx: pip install python-docx")
    doc = Document(filepath)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            # 按标题样式映射为 Markdown 标题
            style = para.style.name.lower() if para.style else ''
            if 'heading 1' in style or '标题 1' in style:
                parts.append(f"# {para.text}")
            elif 'heading 2' in style or '标题 2' in style:
                parts.append(f"## {para.text}")
            elif 'heading 3' in style or '标题 3' in style:
                parts.append(f"### {para.text}")
            else:
                parts.append(para.text)
    # 提取表格
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append('| ' + ' | '.join(cells) + ' |')
    return '\n\n'.join(parts)


def extract_pptx(filepath: str) -> str:
    """从 PowerPoint PPTX 提取文本"""
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("需要安装 python-pptx: pip install python-pptx")
    prs = Presentation(filepath)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"# 幻灯片 {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = ''.join(run.text for run in para.runs).strip()
                    if text:
                        parts.append(text)
        # 备注
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"（备注）{notes}")
    return '\n\n'.join(parts)


def extract_xlsx(filepath: str) -> str:
    """从 Excel XLSX 提取文本"""
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise RuntimeError("需要安装 openpyxl: pip install openpyxl")
    wb = load_workbook(filepath, data_only=True, read_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"## 工作表: {sheet_name}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else '' for c in row]
            if any(cells):
                parts.append('| ' + ' | '.join(cells) + ' |')
    wb.close()
    return '\n\n'.join(parts)


def extract_html(filepath: str) -> str:
    """从 HTML 提取文本"""
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        raise RuntimeError("需要安装 beautifulsoup4: pip install beautifulsoup4")
    with open(filepath, encoding='utf-8', errors='replace') as f:
        soup = BeautifulSoup(f.read(), 'html.parser')
    # 移除脚本和样式
    for tag in soup(['script', 'style', 'nav', 'header', 'footer']):
        tag.decompose()
    parts = []
    for elem in soup.find_all(['h1', 'h2', 'h3', 'h4', 'p', 'li', 'article', 'section']):
        text = elem.get_text(strip=True)
        if not text:
            continue
        if elem.name == 'h1':
            parts.append(f"# {text}")
        elif elem.name == 'h2':
            parts.append(f"## {text}")
        elif elem.name == 'h3':
            parts.append(f"### {text}")
        elif elem.name == 'h4':
            parts.append(f"#### {text}")
        elif elem.name == 'li':
            parts.append(f"- {text}")
        else:
            parts.append(text)
    return '\n\n'.join(parts)


def extract_image(filepath: str, lang='ch_sim+en') -> str:
    """从图片提取文本（OCR）"""
    if not _has_easyocr():
        raise RuntimeError("需要安装 EasyOCR: pip install easyocr")
    import easyocr
    reader = _get_ocr_reader(lang)
    result = reader.readtext(filepath, detail=0, paragraph=True)
    return '\n\n'.join(result)


def ocr_pdf(filepath: str, lang='ch_sim+en', dpi=200) -> str:
    """对图像型 PDF 进行 OCR（先转图片再 OCR）"""
    if not _has_easyocr():
        raise RuntimeError("需要安装 EasyOCR: pip install easyocr")
    try:
        import fitz
    except ImportError:
        raise RuntimeError("需要安装 PyMuPDF: pip install PyMuPDF")
    import easyocr
    reader = _get_ocr_reader(lang)
    doc = fitz.open(filepath)
    parts = []
    for i, page in enumerate(doc, 1):
        # PDF 页面转图片
        mat = fitz.Matrix(dpi / 72, dpi / 72)
        pix = page.get_pixmap(matrix=mat)
        img_path = f"_pdf_page_{os.getpid()}_{i}.png"
        pix.save(img_path)
        try:
            text_lines = reader.readtext(img_path, detail=0, paragraph=True)
            page_text = '\n'.join(text_lines)
        finally:
            if os.path.exists(img_path):
                os.remove(img_path)
        if page_text.strip():
            parts.append(f"\n--- 第 {i} 页 ---\n{page_text}")
    doc.close()
    return '\n\n'.join(parts)


# ============================================================
# 统一调度
# ============================================================
class _EmptyPDFError(Exception):
    pass


_ocr_reader = None


def _get_ocr_reader(lang='ch_sim+en'):
    """懒加载 OCR reader（避免每次重新加载模型）"""
    global _ocr_reader
    if _ocr_reader is None:
        import easyocr
        _ocr_reader = easyocr.Reader(['ch_sim', 'en'], gpu=False)
    return _ocr_reader


def _has_easyocr() -> bool:
    try:
        import easyocr  # noqa
        return True
    except ImportError:
        return False


def extract_any(filepath: str, ocr_lang: str = 'ch_sim+en') -> str:
    """根据文件扩展名自动选择提取方式

    支持：PDF / EPUB / DOCX / PPTX / XLSX / HTML / MD / TXT / 图片
    图像型 PDF 会自动 OCR 回退
    """
    if not os.path.exists(filepath):
        raise FileNotFoundError(f"文件不存在: {filepath}")
    ext = os.path.splitext(filepath)[1].lower()
    base = os.path.basename(filepath)

    if ext == '.pdf':
        try:
            return extract_pdf(filepath)
        except _EmptyPDFError:
            # 图像型 PDF 回退到 OCR
            return ocr_pdf(filepath, lang=ocr_lang)
    elif ext == '.epub':
        return extract_epub(filepath)
    elif ext == '.docx':
        return extract_docx(filepath)
    elif ext == '.pptx':
        return extract_pptx(filepath)
    elif ext == '.xlsx':
        return extract_xlsx(filepath)
    elif ext in ('.html', '.htm'):
        return extract_html(filepath)
    elif ext in ('.md', '.markdown'):
        return read_text(filepath)
    elif ext in ('.txt', '.text'):
        return read_text(filepath)
    elif ext in ('.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.webp', '.gif'):
        return extract_image(filepath, lang=ocr_lang)
    else:
        # 尝试按文本读
        return read_text(filepath)
